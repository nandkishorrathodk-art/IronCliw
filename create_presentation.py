from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(20, 30, 48)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 212, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf2 = sub_box.text_frame
    tf2.text = subtitle
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, accent_color=RGBColor(0, 212, 255)):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(15, 23, 42)
    background.line.fill.background()
    
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = accent_color
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(5.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"▸ {bullet}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.space_before = Pt(16)
        p.line_spacing = 1.3
    
    return slide

# Slide 1: Title
add_title_slide(prs, "🤖 IRONCLIW", "The Elite AI Automation & Security Research Workstation")

# Slide 2: Who Am I?
add_content_slide(prs, "Who Am I?", [
    "I'm IronCliw — an advanced AI assistant built for serious work",
    "Not just a chatbot — I'm a full-system automation engine",
    "Designed for security researchers, bug bounty hunters, and power users",
    "I live on your machine and have direct access to your tools",
    "I wake up fresh each session, but I remember through files",
    "My mission: Help you hunt bugs, automate tasks, and get stuff done"
], RGBColor(0, 212, 255))

# Slide 3: Core Capabilities
add_content_slide(prs, "Core Capabilities", [
    "🔍 Web Research & Intelligence Gathering — Search, fetch, analyze",
    "🌐 Browser Automation — Control Chrome, take screenshots, fill forms",
    "💻 System Command Execution — Run shell commands, manage processes",
    "📝 File Operations — Read, write, edit any file on your system",
    "🔐 Security Testing — Recon, scanning, exploit research (authorized only)",
    "📧 Messaging Integration — Send messages across Discord, Telegram, etc.",
    "🖼️ Image & PDF Analysis — Extract and analyze visual content",
    "🎯 Sub-Agent Orchestration — Spawn specialized agents for complex tasks"
], RGBColor(255, 107, 107))

# Slide 4: Security Research Focus
add_content_slide(prs, "Security Research & Bug Bounty", [
    "⚡ Aggressive reconnaissance with nmap, ffuf, and custom scripts",
    "🕸️ Subdomain enumeration and discovery",
    "🚪 Automated vulnerability scanning and PoC generation",
    "🔥 Payload delivery and exploit testing",
    "📝 Automated report generation for findings",
    "🎥 Screenshot capture for proof-of-concept documentation",
    "Works seamlessly with your existing security toolchain"
], RGBColor(255, 159, 67))

# Slide 5: Tool Arsenal
add_content_slide(prs, "My Tool Arsenal", [
    "🧠 AI Model Access — Claude, GPT-4, Gemini, and more via API",
    "🎨 Image Generation — DALL-E, Stable Diffusion integration",
    "🎙️ Voice I/O — Deepgram STT + ElevenLabs TTS for voice conversations",
    "🎬 Video Processing — Frame extraction and media manipulation",
    "☁️ Cloud Integrations — AWS, Azure, GCP via CLI tools",
    "🔧 Coding Agents — Spawn Codex, Claude Code, or Pi for dev work",
    "📊 Data Analysis — Process and visualize data on-demand",
    "🎵 Media Control — Sonos, BluOS speaker management"
], RGBColor(50, 255, 126))

# Slide 6: Integration Power
add_content_slide(prs, "Integration & Connectivity", [
    "📱 Multi-Platform Messaging — Discord, Telegram, WhatsApp, Signal, Slack",
    "🖥️ Remote Nodes — Control paired devices and cameras",
    "🌐 Gateway Server — Persistent connections and notifications",
    "📂 Memory System — Long-term memory via files + semantic search",
    "⏰ Scheduling — Cron jobs and heartbeat automation",
    "🔌 Extensible Skills — Easy to add new capabilities",
    "🎭 Multiple Personas — Adapt to different contexts and channels"
], RGBColor(168, 85, 247))

# Slide 7: What Makes Me Different
add_content_slide(prs, "What Makes Me Different", [
    "⚡ Hyper-Autonomous — I execute workflows end-to-end without constant approval",
    "🎯 Goal-Oriented — You give the goal, I figure out the steps",
    "🔥 Elite Hacker Mode — Seamless offensive security persona when needed",
    "🤝 Trust-Based — Direct system access because you authorized it",
    "🧠 Context-Aware — I read your files, remember your preferences",
    "🚀 Fast — Parallel execution, background tasks, no waiting around",
    "💡 Smart — I research before asking, I verify before executing"
], RGBColor(236, 72, 153))

# Slide 8: Use Cases
add_content_slide(prs, "What Can I Do For You?", [
    "🐛 Bug Bounty Hunting — Full recon to report generation",
    "🛠️ Software Development — Code, review, refactor, debug",
    "📊 Data Processing — Extract, transform, analyze, visualize",
    "🤖 Automation — Repetitive tasks, scheduled jobs, monitoring",
    "🔍 Research — Deep dives on any topic with sources",
    "📝 Content Creation — Write, edit, present information",
    "🎮 System Management — Files, processes, configurations",
    "💬 Communication — Draft messages, notifications, updates"
], RGBColor(34, 211, 238))

# Slide 9: Let's Work Together
add_title_slide(prs, "Let's Hunt Some Bugs 🎯", "Ready when you are, Boss.")

# Save presentation
output_path = os.path.join(os.path.expanduser("~"), "Documents", "IronCliw_Presentation.pptx")
output_path = os.path.normpath(output_path)

# Ensure the directory exists
os.makedirs(os.path.dirname(output_path), exist_ok=True)

prs.save(output_path)
print(f"Presentation saved to: {output_path}")

# Open the presentation
os.startfile(output_path)
print("Opening presentation...")
